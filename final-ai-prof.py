from fastapi import FastAPI, File, UploadFile, HTTPException, Query, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from haystack.document_stores import ElasticsearchDocumentStore
from haystack.nodes import PreProcessor, FARMReader, ElasticsearchRetriever
from haystack.pipelines import ExtractiveQAPipeline
import uvicorn
import os
import json
import io
from typing import List, Optional
from pydantic import BaseModel
import PyPDF2
from pptx import Presentation
from datetime import datetime

app = FastAPI(title="AI Professor API")

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize document store
document_store = ElasticsearchDocumentStore(
    host="localhost",
    port=9200,
    username="",
    password="",
    index="course_content",
    embedding_dim=384,  # Dimension for the sentence transformers model
    similarity="cosine"
)

# Initialize preprocessor with improved settings
preprocessor = PreProcessor(
    clean_empty_lines=True,
    clean_whitespace=True,
    clean_header_footer=True,
    split_by="word",
    split_length=500,
    split_overlap=50,
    split_respect_sentence_boundary=True
)

# Initialize Reader & Retriever with improved models
reader = FARMReader(
    model_name_or_path="deepset/roberta-base-squad2", 
    use_gpu=True,  # Set to False if no GPU available
    context_window_size=500,
    return_no_answer=True,
    top_k=3
)

retriever = ElasticsearchRetriever(
    document_store=document_store,
    top_k=5
)

# Create pipeline
pipe = ExtractiveQAPipeline(reader, retriever)

class CourseContent(BaseModel):
    title: str
    content: str
    file_type: str
    useful_links: Optional[List[str]] = []

class CourseMetadata(BaseModel):
    course_title: str
    description: Optional[str] = None
    topics: List[str] = []
    date_added: str

@app.post("/upload-content/")
async def upload_content(
    file: UploadFile = File(...),
    course_title: str = Query(..., description="Title of the course"),
    course_description: str = Query(None, description="Course description"),
    topics: str = Query("[]", description="JSON array of course topics"),
    useful_links: str = Query("[]", description="JSON array of useful links")
):
    try:
        # Read file content
        content = await file.read()
        file_type = file.filename.split('.')[-1].lower()
        
        if file_type not in ['pdf', 'pptx']:
            raise HTTPException(
                status_code=400,
                detail="Unsupported file type. Only PDF and PPTX files are allowed."
            )

        text_content = ""

        if file_type == 'pdf':
            # Handle PDF files
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"

        elif file_type == 'pptx':
            # Handle PowerPoint files
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

        # Process links and topics
        try:
            parsed_links = json.loads(useful_links)
            parsed_topics = json.loads(topics)
        except json.JSONDecodeError:
            raise HTTPException(status_code=400, detail="Invalid JSON format in links or topics")

        # Create metadata
        metadata = {
            "course_title": course_title,
            "file_name": file.filename,
            "file_type": file_type,
            "useful_links": parsed_links,
            "topics": parsed_topics,
            "description": course_description,
            "date_added": datetime.now().isoformat(),
            "last_updated": datetime.now().isoformat()
        }

        # Preprocess and index the content
        docs = preprocessor.process([{
            "content": text_content,
            "meta": metadata
        }])
        
        document_store.write_documents(docs)

        return JSONResponse(
            status_code=200,
            content={
                "message": "Content uploaded and indexed successfully",
                "metadata": metadata
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/search/")
async def search(
    question: str = Query(..., description="Question to ask about the course content"),
    course_title: Optional[str] = Query(None, description="Filter by course title")
):
    try:
        # Prepare filters
        filters = {}
        if course_title:
            filters["course_title"] = course_title

        # Get prediction
        prediction = pipe.run(
            query=question,
            params={
                "Retriever": {"top_k": 3, "filters": filters},
                "Reader": {"top_k": 1}
            }
        )

        # Extract answer and context
        if prediction['answers']:
            answer = prediction['answers'][0]
            
            # Get metadata from the document
            metadata = {}
            if prediction['documents']:
                doc = prediction['documents'][0]
                metadata = {
                    "course_title": doc.meta.get('course_title'),
                    "useful_links": doc.meta.get('useful_links', []),
                    "topics": doc.meta.get('topics', []),
                    "file_name": doc.meta.get('file_name')
                }

            return {
                "answer": answer.answer,
                "context": answer.context,
                "confidence": answer.score,
                "metadata": metadata
            }
        else:
            return {
                "answer": "I couldn't find an answer to your question in the course materials.",
                "context": None,
                "confidence": 0.0,
                "metadata": {}
            }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/courses/")
async def list_courses():
    try:
        # Get all documents from the document store
        documents = document_store.get_all_documents()
        
        # Extract course information
        courses = {}
        for doc in documents:
            course_title = doc.meta.get('course_title')
            if course_title not in courses:
                courses[course_title] = {
                    "title": course_title,
                    "description": doc.meta.get('description'),
                    "topics": doc.meta.get('topics', []),
                    "useful_links": doc.meta.get('useful_links', []),
                    "files": [],
                    "date_added": doc.meta.get('date_added')
                }
            courses[course_title]["files"].append({
                "name": doc.meta.get('file_name'),
                "type": doc.meta.get('file_type'),
                "last_updated": doc.meta.get('last_updated')
            })
        
        return {"courses": list(courses.values())}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/health/")
async def health_check():
    try:
        # Check Elasticsearch connection
        es_status = document_store.client.ping()
        
        # Check if the reader model is loaded
        reader_status = reader.model is not None
        
        return {
            "status": "healthy" if es_status and reader_status else "unhealthy",
            "elasticsearch": "connected" if es_status else "disconnected",
            "model": "loaded" if reader_status else "not loaded"
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "error": str(e)
        }

if __name__ == "__main__":
    uvicorn.run("main.py:app", host="0.0.0.0", port=8000, reload=True)
